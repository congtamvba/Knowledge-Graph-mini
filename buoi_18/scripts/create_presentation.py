from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "outputs" / "Buoi_18_AI_Compliance_Audit_polished.pptx"

NAVY = RGBColor(22, 42, 58)
INK = RGBColor(30, 42, 46)
MUTED = RGBColor(91, 105, 104)
PAPER = RGBColor(247, 246, 241)
WHITE = RGBColor(255, 255, 255)
TEAL = RGBColor(26, 119, 117)
MINT = RGBColor(218, 239, 231)
GOLD = RGBColor(224, 166, 67)
GOLD_SOFT = RGBColor(250, 239, 207)
CORAL = RGBColor(204, 87, 75)
CORAL_SOFT = RGBColor(250, 225, 220)
LINE = RGBColor(215, 224, 218)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


def rect(slide, x, y, w, h, fill, radius=False, line=None):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    if radius:
        shape.adjustments[0] = 0.12
    return shape


def text(slide, value, x, y, w, h, size=18, color=INK, bold=False, font="Aptos", align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.04)
    frame.margin_right = Inches(0.04)
    frame.margin_top = Inches(0.02)
    frame.vertical_anchor = valign
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = value
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def bullet_list(slide, items, x, y, w, h, size=16, color=INK):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.08)
    frame.margin_right = Inches(0.04)
    for index, item in enumerate(items):
        p = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(10)
        p.bullet = True
    return box


def base(slide, number, title, kicker="BUỔI 18 / AI COMPLIANCE & AUDIT"):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PAPER
    rect(slide, 0, 0, 13.333, 0.16, TEAL)
    text(slide, kicker, 0.55, 0.42, 6.2, 0.25, 10, TEAL, True)
    text(slide, title, 0.55, 0.78, 10.9, 0.62, 27, NAVY, True, "Georgia")
    rect(slide, 0.55, 1.48, 1.15, 0.06, GOLD)
    rect(slide, 11.96, 0.38, 0.72, 0.5, NAVY, True)
    text(slide, f"{number:02d}", 11.96, 0.49, 0.72, 0.22, 12, WHITE, True, align=PP_ALIGN.CENTER)
    rect(slide, 0.55, 6.98, 12.2, 0.012, LINE)
    text(slide, "Kết quả AI là bản nháp hỗ trợ kiểm toán viên, không thay thế quyết định con người.", 0.55, 7.08, 8, 0.2, 9, MUTED)


def pill(slide, label, x, y, w, fill, color=INK):
    rect(slide, x, y, w, 0.36, fill, True)
    text(slide, label, x, y + 0.05, w, 0.22, 10, color, True, align=PP_ALIGN.CENTER)


def arrow(slide, x, y, w, color=TEAL):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(0.32))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color

# 1
slide = prs.slides.add_slide(blank)
slide.background.fill.solid(); slide.background.fill.fore_color.rgb = NAVY
rect(slide, 0, 0, 13.333, 0.18, GOLD)
rect(slide, 0.62, 0.7, 0.08, 5.75, GOLD)
text(slide, "BUỔI 18", 1.0, 0.78, 2.2, 0.35, 15, GOLD, True)
text(slide, "AI Compliance\n& Audit", 1.0, 1.35, 6.5, 1.55, 38, WHITE, True, "Georgia")
text(slide, "Từ quy định phân tán đến bằng chứng kiểm toán có kiểm soát", 1.05, 3.35, 5.7, 0.55, 19, RGBColor(220, 232, 229))
pill(slide, "UC3  COMPLIANCE CHECKER", 1.05, 5.02, 2.55, TEAL, WHITE)
pill(slide, "UC4  AUDIT CHECKLIST", 3.78, 5.02, 2.35, GOLD, NAVY)
text(slide, "Agribank internal policies  ·  legal evidence  ·  RBAC  ·  audit trail", 1.05, 5.75, 8.5, 0.3, 12, RGBColor(190, 210, 207))
# self-contained hero infographic: source -> evidence -> accountable review
rect(slide, 8.15, 1.0, 4.25, 4.95, RGBColor(31, 59, 73), True)
text(slide, "GOVERNED\nDECISION FLOW", 8.55, 1.45, 3.25, 0.75, 18, WHITE, True, "Georgia")
for i, (label, sub, fill) in enumerate([("SOURCE", "policy + law", MINT), ("EVIDENCE", "citation + article", GOLD_SOFT), ("REVIEW", "human in the loop", CORAL_SOFT)]):
    y = 2.55 + i * 1.0
    rect(slide, 8.58, y, 3.35, 0.65, fill, True)
    text(slide, label, 8.78, y + 0.1, 1.2, 0.22, 12, NAVY, True)
    text(slide, sub, 10.0, y + 0.1, 1.7, 0.22, 12, MUTED, True)
    if i < 2:
        arrow(slide, 10.0, y + 0.73, 0.45, GOLD)

# 2
slide = prs.slides.add_slide(blank); base(slide, 2, "Bài toán: tìm đúng quy định, đúng quyền, đúng bằng chứng")
for x, title, body, fill in [(0.7, "01 · PHÂN TÁN", "Nội bộ + pháp luật\nở nhiều nguồn", CORAL_SOFT), (4.55, "02 · RỦI RO", "Mâu thuẫn hạn mức,\nquy trình, thẩm quyền", GOLD_SOFT), (8.4, "03 · KIỂM SOÁT", "Citation thật + RBAC\n+ human review", MINT)]:
    rect(slide, x, 1.85, 3.45, 2.2, fill, True)
    text(slide, title, x + 0.22, 2.15, 2.9, 0.3, 13, NAVY, True)
    text(slide, body, x + 0.22, 2.72, 2.95, 0.8, 22, INK, True, "Georgia")
    rect(slide, x + 0.22, 3.62, 2.0, 0.04, TEAL)
arrow(slide, 4.0, 2.8, 0.42, TEAL); arrow(slide, 7.85, 2.8, 0.42, TEAL)
text(slide, "Mục tiêu: biến truy xuất quy định thành một quy trình có thể kiểm chứng, truy vết và bàn giao cho kiểm toán viên.", 0.85, 5.05, 11.4, 0.62, 21, NAVY, True, "Georgia", align=PP_ALIGN.CENTER)

# 3
slide = prs.slides.add_slide(blank); base(slide, 3, "Dữ liệu nền: đủ rộng để đối chiếu, đủ chặt để kiểm soát")
for x, value, label, fill in [(0.8, "24", "chunks nội bộ", MINT), (3.75, "811", "chunks kết hợp", GOLD_SOFT), (6.7, "10", "domain nội bộ", CORAL_SOFT), (9.65, "14/14", "metadata đầy đủ", MINT)]:
    rect(slide, x, 1.75, 2.25, 1.62, fill, True)
    text(slide, value, x, 2.02, 2.25, 0.6, 31, NAVY, True, "Georgia", align=PP_ALIGN.CENTER)
    text(slide, label, x + 0.15, 2.85, 1.95, 0.28, 12, MUTED, True, align=PP_ALIGN.CENTER)
text(slide, "Các trường kiểm soát quan trọng", 0.85, 4.12, 3.2, 0.3, 15, TEAL, True)
for i, (label, value) in enumerate([("article", "811/811"), ("citation", "811/811"), ("allowed_roles", "811/811")]):
    y = 4.65 + i * 0.52
    text(slide, label, 0.95, y, 2.1, 0.25, 15, INK, True)
    rect(slide, 3.0, y + 0.05, 5.5, 0.16, LINE, True)
    rect(slide, 3.0, y + 0.05, 5.5, 0.16, TEAL, True)
    text(slide, value, 8.8, y - 0.03, 1.2, 0.25, 13, TEAL, True)
text(slide, "Nguồn được đọc read-only; citation và quyền truy cập nằm trong metadata thật.", 8.7, 4.65, 3.3, 1.1, 20, NAVY, True, "Georgia")

# 4
slide = prs.slides.add_slide(blank); base(slide, 4, "Kiến trúc kiểm soát: lọc quyền trước khi tìm kiếm")
steps = [("NGƯỜI DÙNG", "role + request", NAVY), ("RBAC", "allowed_roles", TEAL), ("BM25", "evidence ranking", GOLD), ("LLM / FALLBACK", "JSON guardrail", CORAL)]
for i, (head, body, fill) in enumerate(steps):
    x = 0.75 + i * 3.05
    rect(slide, x, 2.05, 2.4, 1.45, fill, True)
    text(slide, head, x + 0.12, 2.35, 2.16, 0.27, 14, WHITE, True, align=PP_ALIGN.CENTER)
    text(slide, body, x + 0.12, 2.85, 2.16, 0.28, 13, WHITE, False, align=PP_ALIGN.CENTER)
    if i < 3: arrow(slide, x + 2.47, 2.62, 0.42, GOLD)
text(slide, "Evidence package", 0.9, 4.28, 2.0, 0.28, 14, TEAL, True)
rect(slide, 0.9, 4.72, 11.55, 0.85, WHITE, True, LINE)
for i, label in enumerate(["CITATION", "ARTICLE", "TEXT", "REVIEW_STATUS"]):
    pill(slide, label, 1.2 + i * 2.7, 4.97, 1.9, MINT if i != 3 else GOLD_SOFT, TEAL if i != 3 else NAVY)
text(slide, "Không có chunk được phép truy cập thì không được đưa vào retrieval hoặc context.", 1.0, 6.03, 11.2, 0.3, 16, NAVY, True, align=PP_ALIGN.CENTER)

# 5
slide = prs.slides.add_slide(blank); base(slide, 5, "UC3 · Compliance Checker: đối chiếu hai phía của quy định")
text(slide, "Văn bản nội bộ", 0.8, 1.7, 2.5, 0.3, 16, TEAL, True)
rect(slide, 0.8, 2.15, 4.45, 2.15, MINT, True)
text(slide, "100/QĐ-NHNO-AT\n250/QĐ-NHNO-QLRR\n315/QC-NHNO-TD", 1.15, 2.55, 3.7, 0.9, 22, NAVY, True, "Georgia")
text(slide, "hạn mức · phương tiện · thẩm quyền", 1.15, 3.72, 3.7, 0.25, 13, MUTED)
arrow(slide, 5.55, 2.85, 1.0, GOLD)
text(slide, "Evidence\npackage", 5.55, 3.45, 1.0, 0.5, 13, NAVY, True, align=PP_ALIGN.CENTER)
text(slide, "Văn bản pháp luật", 7.2, 1.7, 2.8, 0.3, 16, CORAL, True)
rect(slide, 7.2, 2.15, 5.2, 2.15, CORAL_SOFT, True)
text(slide, "01/2014/TT-NHNN\n41/2016/TT-NHNN\nNghị định / Luật liên quan", 7.55, 2.55, 4.45, 0.9, 21, NAVY, True, "Georgia")
text(slide, "điều · khoản · hiệu lực · phạm vi", 7.55, 3.72, 4.45, 0.25, 13, MUTED)
rect(slide, 1.05, 5.15, 11.2, 0.7, WHITE, True, LINE)
text(slide, "Kết quả P2", 1.35, 5.36, 1.5, 0.25, 15, TEAL, True)
text(slide, "3 findings  ·  Citation PASS  ·  Human review REQUIRED", 3.0, 5.36, 8.5, 0.25, 17, NAVY, True)

# 6
slide = prs.slides.add_slide(blank); base(slide, 6, "UC4 · Audit Checklist: từ điều khoản đến câu hỏi kiểm toán")
text(slide, "An toàn kho quỹ", 0.8, 1.65, 3.2, 0.3, 17, TEAL, True)
text(slide, "Bảo mật CNTT & AI", 6.9, 1.65, 3.4, 0.3, 17, CORAL, True)
for x, items, fill in [(0.8, [("CHK_01", "Xe bọc thép / phương án bảo vệ", "HIGH"), ("CHK_02", "Kiểm đếm / niêm phong", "MEDIUM"), ("CHK_03", "Ban quản lý kho tiền", "MEDIUM")], MINT), (6.9, [("CHK_01", "Mã hóa dữ liệu AI", "HIGH"), ("CHK_02", "Lưu audit trail 12 tháng", "HIGH")], CORAL_SOFT)]:
    rect(slide, x, 2.1, 5.55, 3.1, fill, True)
    for i, (code, label, risk) in enumerate(items):
        y = 2.42 + i * 0.8
        text(slide, code, x + 0.28, y, 0.85, 0.25, 12, NAVY, True)
        text(slide, label, x + 1.25, y, 3.35, 0.3, 15, INK, True)
        pill(slide, risk, x + 4.65, y - 0.03, 0.65, GOLD_SOFT if risk == "HIGH" else MINT, NAVY)
text(slide, "Kết quả P3: 6 items · citations attached · mọi item cần xác minh", 1.0, 5.78, 11.2, 0.32, 18, NAVY, True, "Georgia", align=PP_ALIGN.CENTER)

# 7
slide = prs.slides.add_slide(blank); base(slide, 7, "Guardrails: AI có thể gợi ý, nhưng không tự phê duyệt")
checks = [("RBAC", "Staff không thấy dữ liệu restricted", MINT), ("CITATION", "Citation phải tồn tại trong corpus", GOLD_SOFT), ("PRIVACY", "Log không chứa secret", CORAL_SOFT), ("REVIEW", "NEEDS_HUMAN_REVIEW luôn bật", MINT)]
for i, (head, body, fill) in enumerate(checks):
    x = 0.85 + (i % 2) * 6.05
    y = 1.75 + (i // 2) * 2.15
    rect(slide, x, y, 5.35, 1.55, fill, True)
    text(slide, "✓", x + 0.25, y + 0.28, 0.5, 0.55, 30, TEAL, True)
    text(slide, head, x + 0.95, y + 0.3, 2.2, 0.27, 15, NAVY, True)
    text(slide, body, x + 0.95, y + 0.78, 3.9, 0.35, 15, INK)
text(slide, "P5 SECURITY & GUARDRAIL TESTS: 7/7 PASS", 1.0, 6.05, 11.2, 0.35, 20, TEAL, True, "Georgia", align=PP_ALIGN.CENTER)

# 8
slide = prs.slides.add_slide(blank); base(slide, 8, "Streamlit workbench: ba luồng trong một màn hình")
for x, label, detail, fill in [(0.8, "UC3", "Compliance\nChecker", MINT), (4.35, "UC4", "Audit\nChecklist", GOLD_SOFT), (7.9, "AUDIT", "System\nTrail", CORAL_SOFT)]:
    rect(slide, x, 1.85, 2.85, 2.2, fill, True)
    text(slide, label, x + 0.25, 2.16, 1.0, 0.3, 18, TEAL if label == "UC3" else CORAL if label == "AUDIT" else NAVY, True)
    text(slide, detail, x + 0.25, 2.78, 2.25, 0.75, 25, NAVY, True, "Georgia")
    pill(slide, "EXPORT", x + 0.25, 3.62, 0.9, WHITE, TEAL)
arrow(slide, 3.78, 2.82, 0.42, TEAL); arrow(slide, 7.32, 2.82, 0.42, TEAL)
rect(slide, 1.15, 4.75, 10.9, 0.9, NAVY, True)
text(slide, "Banner cảnh báo", 1.45, 5.02, 2.1, 0.25, 14, GOLD, True)
text(slide, "Kết quả gợi ý cần kiểm toán viên xác minh trước khi ban hành.", 3.65, 5.0, 7.8, 0.28, 16, WHITE, True)

# 9
slide = prs.slides.add_slide(blank); base(slide, 9, "Kết quả kiểm thử: hệ thống đã sẵn sàng để trình diễn")
metrics = [("P1", "10 domains", "CATALOG PASS", TEAL), ("P2", "3 findings", "CHECKER PASS", GOLD), ("P3", "6 checklist", "GENERATOR PASS", CORAL), ("P5", "7 / 7", "SECURITY PASS", TEAL), ("P6", "8 / 8", "VALIDATION PASS", NAVY)]
for i, (stage, value, status, color) in enumerate(metrics):
    x = 0.75 + i * 2.48
    rect(slide, x, 1.85, 2.1, 2.0, WHITE, True, LINE)
    rect(slide, x, 1.85, 2.1, 0.16, color)
    text(slide, stage, x + 0.2, 2.18, 0.5, 0.25, 13, color, True)
    text(slide, value, x + 0.2, 2.7, 1.7, 0.45, 23, NAVY, True, "Georgia")
    text(slide, status, x + 0.2, 3.42, 1.7, 0.24, 10, MUTED, True)
rect(slide, 1.0, 4.72, 11.3, 1.0, MINT, True)
text(slide, "SYSTEM READY FOR DEMO: YES", 1.25, 5.02, 10.8, 0.35, 25, TEAL, True, "Georgia", align=PP_ALIGN.CENTER)

# 10
slide = prs.slides.add_slide(blank); base(slide, 10, "Thông điệp mang về")
text(slide, "Một hệ thống compliance tốt không chỉ trả lời nhanh.", 0.9, 1.65, 11.3, 0.5, 25, NAVY, True, "Georgia", align=PP_ALIGN.CENTER)
text(slide, "Nó phải trả lời từ đúng nguồn, đúng quyền,\nđúng citation và đúng người chịu trách nhiệm cuối cùng.", 1.25, 2.45, 10.7, 1.0, 28, TEAL, True, "Georgia", align=PP_ALIGN.CENTER)
for i, (label, fill) in enumerate([("SOURCE", MINT), ("ACCESS", GOLD_SOFT), ("EVIDENCE", CORAL_SOFT), ("REVIEW", MINT)]):
    x = 1.15 + i * 3.0
    rect(slide, x, 4.35, 2.35, 0.9, fill, True)
    text(slide, label, x, 4.65, 2.35, 0.3, 15, NAVY, True, align=PP_ALIGN.CENTER)
    if i < 3: arrow(slide, x + 2.42, 4.63, 0.35, GOLD)
text(slide, "Cảm ơn", 5.2, 6.0, 2.9, 0.45, 25, GOLD, True, "Georgia", align=PP_ALIGN.CENTER)

OUTPUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(OUTPUT)
print(f"PRESENTATION={OUTPUT}")
print(f"SLIDES={len(prs.slides)}")
